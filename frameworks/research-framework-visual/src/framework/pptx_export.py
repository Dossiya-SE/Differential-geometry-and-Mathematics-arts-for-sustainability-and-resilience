from __future__ import annotations
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN


def _rgb(hexcolor):
    h=hexcolor.lstrip('#'); return RGBColor(int(h[0:2],16),int(h[2:4],16),int(h[4:6],16))


def export_pptx(request, framework_spec, path):
    prs=Presentation(); prs.slide_width=Inches(15); prs.slide_height=Inches(10)
    slide=prs.slides.add_slide(prs.slide_layouts[6])
    bg=slide.background.fill; bg.solid(); bg.fore_color.rgb=_rgb('#081018')
    title=slide.shapes.add_textbox(Inches(.3),Inches(.15),Inches(14.4),Inches(.45))
    p=title.text_frame.paragraphs[0]; p.text=request['canvas']['title']; p.font.size=Pt(20); p.font.bold=True; p.font.color.rgb=_rgb('#ffffff'); p.alignment=PP_ALIGN.CENTER
    sub=slide.shapes.add_textbox(Inches(.3),Inches(.58),Inches(14.4),Inches(.28))
    p=sub.text_frame.paragraphs[0]; p.text=request['canvas']['subtitle']; p.font.size=Pt(11); p.font.italic=True; p.font.color.rgb=_rgb('#cfd8e3'); p.alignment=PP_ALIGN.CENTER
    stages=framework_spec['stages']; gap=.05
    top_w=(14.8-gap*9)/10
    for i,s in enumerate(stages[:10]):
        x=.1+i*(top_w+gap); y=.9; h=2.15
        sh=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(top_w),Inches(h)); sh.fill.solid(); sh.fill.fore_color.rgb=_rgb('#0c1322'); sh.line.color.rgb=_rgb(s['color']); sh.line.width=Pt(1.5)
        tb=slide.shapes.add_textbox(Inches(x+.05),Inches(y+.08),Inches(top_w-.1),Inches(.55)); tf=tb.text_frame; tf.clear(); p=tf.paragraphs[0]; p.text=f"{s['id']}  {s['title']}"; p.font.size=Pt(8); p.font.bold=True; p.font.color.rgb=_rgb(s['color']); p.alignment=PP_ALIGN.CENTER
        p=tf.add_paragraph(); p.text=s['body']; p.font.size=Pt(6.5); p.font.color.rgb=_rgb('#ffffff'); p.alignment=PP_ALIGN.CENTER
    central=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(.1),Inches(3.2),Inches(14.8),Inches(3.35)); central.fill.solid(); central.fill.fore_color.rgb=_rgb('#050b14'); central.line.color.rgb=_rgb('#7a4fd6')
    tb=slide.shapes.add_textbox(Inches(.5),Inches(3.45),Inches(14),Inches(.55)); p=tb.text_frame.paragraphs[0]; p.text='Evidence field → Gap space → Research question → Model manifold → Verification → Calibration → Validation'; p.font.size=Pt(13); p.font.color.rgb=_rgb('#42d4e5'); p.alignment=PP_ALIGN.CENTER
    tb=slide.shapes.add_textbox(Inches(4.1),Inches(4.35),Inches(6.8),Inches(.7)); p=tb.text_frame.paragraphs[0]; p.text='ẋ = f(x,u,η,θ,t)\n g(x,u) ≤ 0'; p.font.size=Pt(16); p.font.italic=True; p.font.color.rgb=_rgb('#ffffff'); p.alignment=PP_ALIGN.CENTER
    bottom_w=(14.8-gap*5)/6
    for i,s in enumerate(stages[10:]):
        x=.1+i*(bottom_w+gap); y=6.7; h=2.2
        sh=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(bottom_w),Inches(h)); sh.fill.solid(); sh.fill.fore_color.rgb=_rgb('#0c1322'); sh.line.color.rgb=_rgb(s['color']); sh.line.width=Pt(1.5)
        tb=slide.shapes.add_textbox(Inches(x+.05),Inches(y+.08),Inches(bottom_w-.1),Inches(.7)); tf=tb.text_frame; tf.clear(); p=tf.paragraphs[0]; p.text=f"{s['id']}  {s['title']}"; p.font.size=Pt(9); p.font.bold=True; p.font.color.rgb=_rgb(s['color']); p.alignment=PP_ALIGN.CENTER
        p=tf.add_paragraph(); p.text=s['body']; p.font.size=Pt(7); p.font.color.rgb=_rgb('#ffffff'); p.alignment=PP_ALIGN.CENTER
    note=slide.shapes.add_textbox(Inches(.4),Inches(9.15),Inches(14.2),Inches(.5)); p=note.text_frame.paragraphs[0]; p.text='Convenience-editable derivative. Canonical fidelity authority: poster_EDITABLE.svg'; p.font.size=Pt(8); p.font.color.rgb=_rgb('#cfd8e3'); p.alignment=PP_ALIGN.CENTER
    Path(path).parent.mkdir(parents=True,exist_ok=True); prs.save(path)
