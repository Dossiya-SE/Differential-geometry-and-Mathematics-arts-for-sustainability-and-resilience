from __future__ import annotations

import importlib
import importlib.metadata as metadata
import os
import platform
import shutil
import sys
from pathlib import Path


def _inside(path: str | Path, parent: str | Path) -> bool:
    try:
        Path(path).resolve().relative_to(Path(parent).resolve())
        return True
    except ValueError:
        return False


def _version(distribution: str) -> str | None:
    try:
        return metadata.version(distribution)
    except metadata.PackageNotFoundError:
        return None


def _import_check(module: str) -> bool:
    try:
        importlib.import_module(module)
        return True
    except Exception:
        return False


def doctor(require_notebook: bool = False) -> dict:
    """Audit interpreter isolation and rendering dependencies.

    The report is intentionally machine-readable so GitHub Actions and local
    MacBook setup can use exactly the same checks.
    """
    executable = Path(sys.executable).resolve()
    prefix = Path(sys.prefix).resolve()
    venv = os.environ.get("VIRTUAL_ENV")
    conda = os.environ.get("CONDA_PREFIX")
    jupyter_executable = shutil.which("jupyter")

    checks: dict[str, bool] = {}
    warnings: list[str] = []
    failures: list[str] = []

    checks["python_supported"] = sys.version_info >= (3, 10)
    if not checks["python_supported"]:
        failures.append("Python >= 3.10 is required.")

    if venv:
        checks["venv_interpreter_isolated"] = _inside(executable, venv)
        if not checks["venv_interpreter_isolated"]:
            failures.append(
                f"VIRTUAL_ENV={venv} but sys.executable={executable}; interpreter leakage detected."
            )
    else:
        checks["venv_interpreter_isolated"] = True

    if venv and conda and Path(venv).resolve() != Path(conda).resolve():
        warnings.append(
            "A Conda environment and a venv are both visible. Use the venv's exact python path for Jupyter and rendering."
        )

    core = {
        "CairoSVG": "cairosvg",
        "Pillow": "PIL",
        "PyYAML": "yaml",
        "python-pptx": "pptx",
    }
    for dist, module in core.items():
        ok = _import_check(module)
        checks[f"import_{module}"] = ok
        if not ok:
            failures.append(f"Required dependency cannot be imported: {dist} ({module}).")

    # Real conversion smoke test: catches a Python package that imports while
    # its native Cairo runtime is unusable.
    try:
        import cairosvg

        png = cairosvg.svg2png(
            bytestring=b'<svg xmlns="http://www.w3.org/2000/svg" width="8" height="8"><rect width="8" height="8" fill="black"/></svg>'
        )
        checks["cairosvg_conversion"] = png.startswith(b"\x89PNG")
    except Exception as exc:
        checks["cairosvg_conversion"] = False
        failures.append(f"CairoSVG conversion smoke test failed: {exc}")

    try:
        from pptx import Presentation

        presentation = Presentation()
        presentation.slides.add_slide(presentation.slide_layouts[6])
        checks["pptx_creation"] = len(presentation.slides) == 1
    except Exception as exc:
        checks["pptx_creation"] = False
        failures.append(f"PPTX creation smoke test failed: {exc}")

    notebook_versions: dict[str, str | None] = {}
    if require_notebook:
        notebook_modules = {
            "jupyter": "jupyter",
            "notebook": "notebook",
            "ipykernel": "ipykernel",
            "numpy": "numpy",
            "matplotlib": "matplotlib",
        }
        for dist, module in notebook_modules.items():
            ok = _import_check(module)
            checks[f"notebook_import_{module}"] = ok
            notebook_versions[dist] = _version(dist)
            if not ok:
                failures.append(f"Notebook dependency cannot be imported: {dist}.")

        if jupyter_executable:
            checks["jupyter_same_environment"] = _inside(jupyter_executable, prefix)
            if not checks["jupyter_same_environment"]:
                failures.append(
                    "Jupyter executable is outside the active Python environment: "
                    f"{jupyter_executable}. Launch with '<venv>/bin/python -m notebook'."
                )
        else:
            checks["jupyter_same_environment"] = False
            failures.append("No Jupyter executable is available in the active environment.")

        try:
            import matplotlib
            from PIL import ImageFont

            font_dir = Path(matplotlib.get_data_path()) / "fonts" / "ttf"
            required_fonts = [
                "DejaVuSans.ttf",
                "DejaVuSans-Bold.ttf",
                "DejaVuSerif.ttf",
                "DejaVuSerif-Bold.ttf",
                "DejaVuSerif-Italic.ttf",
            ]
            for name in required_fonts:
                path = font_dir / name
                if not path.exists():
                    raise FileNotFoundError(path)
                ImageFont.truetype(str(path), 18)
            checks["matplotlib_bundled_fonts"] = True
        except Exception as exc:
            checks["matplotlib_bundled_fonts"] = False
            failures.append(f"Matplotlib bundled-font check failed: {exc}")

    versions = {
        "CairoSVG": _version("CairoSVG"),
        "Pillow": _version("Pillow"),
        "PyYAML": _version("PyYAML"),
        "python-pptx": _version("python-pptx"),
        **notebook_versions,
    }

    return {
        "pass": not failures and all(checks.values()),
        "platform": platform.system(),
        "machine": platform.machine(),
        "python": platform.python_version(),
        "sys_executable": str(executable),
        "sys_prefix": str(prefix),
        "virtual_env": venv,
        "conda_prefix": conda,
        "jupyter_executable": jupyter_executable,
        "versions": versions,
        "checks": checks,
        "warnings": warnings,
        "failures": failures,
    }
